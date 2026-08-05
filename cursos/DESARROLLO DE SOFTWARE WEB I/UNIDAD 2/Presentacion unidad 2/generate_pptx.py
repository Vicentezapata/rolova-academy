#!/usr/bin/env python3
"""
Generador de PPTX para Unidad 2 — IF203IINF
Usa python-pptx para crear una presentación con los slides del workflow PPT Agent
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path

# === Configuración ===
OUTPUT_DIR = Path("ppt-output")
SLIDES_DIR = OUTPUT_DIR / "slides"
IMAGES_DIR = OUTPUT_DIR / "images"
OUTPUT_PPTX = OUTPUT_DIR / "presentacion_U2_new.pptx"

# Dimensiones 16:9 (1280x720 equivalente en EMU)
SLIDE_WIDTH = Inches(13.33)   # 1280px ≈ 13.33in @ 96dpi
SLIDE_HEIGHT = Inches(7.5)    # 720px ≈ 7.5in

# === Paleta dark_tech ===
BG = RGBColor(0x05, 0x0b, 0x1f)
BG2 = RGBColor(0x0a, 0x1f, 0x3d)
ACCENT1 = RGBColor(0x22, 0xD3, 0xEE)  # cyan
ACCENT2 = RGBColor(0x63, 0x66, 0xf1)  # indigo
ACCENT3 = RGBColor(0xFD, 0xE0, 0x47)  # yellow
ACCENT4 = RGBColor(0x10, 0xb9, 0x81)  # green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIM = RGBColor(0xCB, 0xD5, 0xE1)
ERR = RGBColor(0xEF, 0x44, 0x44)

def set_slide_bg(slide, color: RGBColor):
    """Relleno de fondo sólido"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=WHITE,
                 font_name="Inter", align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def make_slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG)
    W, H = SLIDE_WIDTH, SLIDE_HEIGHT
    
    # Try to add cover image
    img_path = IMAGES_DIR / "cover.png"
    if img_path.exists():
        pic = slide.shapes.add_picture(str(img_path), Inches(6), 0, Inches(7.33), H)
    
    # Gradient overlay rect (simulate)
    add_rect(slide, 0, 0, Inches(9), H, fill_color=BG)
    
    # Tag line
    add_text_box(slide, "IF203IINF  ·  IPSS  ·  2026", Inches(0.8), Inches(1.2), Inches(6), Inches(0.4),
                 font_size=10, color=ACCENT1, font_name="Calibri")
    
    # Title
    add_text_box(slide, "UNIDAD 2", Inches(0.8), Inches(1.8), Inches(7), Inches(0.6),
                 font_size=14, bold=True, color=ACCENT1, font_name="Calibri")
    add_text_box(slide, "Planificación de Pruebas", Inches(0.8), Inches(2.5), Inches(8), Inches(1.5),
                 font_size=40, bold=True, color=WHITE, font_name="Calibri")
    
    # Subtitle
    add_text_box(slide, 
        "Estrategia, diseño y documentación de pruebas de software.\nDesde la selección contextual hasta el Plan Maestro IEEE 829.",
        Inches(0.8), Inches(4.2), Inches(7.5), Inches(1),
        font_size=14, color=DIM, font_name="Calibri")
    
    # Bottom chips
    for i, (lbl, txt) in enumerate([("4 Sesiones", ACCENT2), ("18 Slides", ACCENT2), ("Gherkin · RTM · Test Plan", DIM)]):
        box = add_rect(slide, Inches(0.8 + i * 2.5), Inches(5.5), Inches(2.3), Inches(0.4), line_color=ACCENT2)
        add_text_box(slide, lbl, Inches(0.85 + i * 2.5), Inches(5.52), Inches(2.2), Inches(0.38),
                     font_size=11, color=WHITE, font_name="Calibri", align=PP_ALIGN.CENTER)

def make_toc_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    W, H = SLIDE_WIDTH, SLIDE_HEIGHT
    
    add_text_box(slide, "UNIDAD 2 — IF203IINF", Inches(0.5), Inches(0.3), Inches(5), Inches(0.4),
                 font_size=10, color=ACCENT1, font_name="Calibri")
    add_text_box(slide, "Contenido de la Unidad", Inches(0.5), Inches(0.7), Inches(8), Inches(0.5),
                 font_size=22, bold=True, color=WHITE, font_name="Calibri")
    # Divider
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.02), fill_color=ACCENT1)
    
    sessions = [
        ("Sesión 01", "Selección Estratégica de Tipos de Prueba",
         "Contexto ISTQB · Matriz Riesgo Impacto×Probabilidad · MVP vs App Crítica", ACCENT1),
        ("Sesión 02", "Requerimientos e Historias de Usuario",
         "User Stories INVEST · Sintaxis Gherkin BDD · DADO/CUANDO/ENTONCES", ACCENT2),
        ("Sesión 03", "Diseño de Casos de Prueba y Trazabilidad",
         "Anatomía Test Case (5 campos) · RTM · Cobertura · Auditoría", ACCENT3),
        ("Sesión 04 — Sincrónica", "El Plan de Pruebas Maestro",
         "IEEE 829 · Alcance/Estrategia/Criterios · Roleplay + QA Audit", ACCENT4),
    ]
    
    cols = 2
    card_w = Inches(5.9)
    card_h = Inches(2.4)
    gap = Inches(0.25)
    start_x = Inches(0.5)
    start_y = Inches(1.4)
    
    for i, (session, title, topics, color) in enumerate(sessions):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        
        # Card bg
        add_rect(slide, x, y, card_w, card_h, line_color=color)
        # Session label
        add_text_box(slide, session, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.3), Inches(0.3),
                     font_size=10, color=color, font_name="Calibri")
        # Title
        add_text_box(slide, title, x + Inches(0.2), y + Inches(0.5), card_w - Inches(0.3), Inches(0.8),
                     font_size=14, bold=True, color=WHITE, font_name="Calibri")
        # Topics
        add_text_box(slide, topics, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.3), Inches(0.8),
                     font_size=11, color=DIM, font_name="Calibri")

def make_section_slide(prs, part_label, title, subtitle, accent_color, agenda_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    
    # Center content
    add_text_box(slide, part_label, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.4),
                 font_size=12, color=accent_color, font_name="Calibri", align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5.5), Inches(2.3), Inches(2.3), Inches(0.02), fill_color=accent_color)
    
    add_text_box(slide, title, Inches(1), Inches(2.4), Inches(11.3), Inches(1.2),
                 font_size=36, bold=True, color=WHITE, font_name="Calibri", align=PP_ALIGN.CENTER)
    add_text_box(slide, subtitle, Inches(2), Inches(3.7), Inches(9.3), Inches(0.8),
                 font_size=15, color=DIM, font_name="Calibri", align=PP_ALIGN.CENTER)
    
    # Agenda chips
    chip_w = Inches(3.5)
    chip_h = Inches(0.8)
    total_w = len(agenda_items) * chip_w + (len(agenda_items) - 1) * Inches(0.2)
    start_x = (SLIDE_WIDTH - total_w) / 2
    
    for i, (label, text) in enumerate(agenda_items):
        x = start_x + i * (chip_w + Inches(0.2))
        add_rect(slide, x, Inches(4.8), chip_w, chip_h, line_color=accent_color)
        add_text_box(slide, label, x + Inches(0.1), Inches(4.85), chip_w - Inches(0.2), Inches(0.3),
                     font_size=9, color=accent_color, font_name="Calibri", align=PP_ALIGN.CENTER)
        add_text_box(slide, text, x + Inches(0.1), Inches(5.15), chip_w - Inches(0.2), Inches(0.4),
                     font_size=11, color=DIM, font_name="Calibri", align=PP_ALIGN.CENTER)

def make_content_slide(prs, session_tag, title, content_blocks, accent_color=ACCENT1, img_path=None):
    """Generic content slide with 2-3 content blocks"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    
    # If image provided, add on right side
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), Inches(8.5), Inches(1), Inches(4.5), Inches(6))
        add_rect(slide, Inches(8), 0, Inches(1), SLIDE_HEIGHT, fill_color=BG)
    
    # Header
    add_text_box(slide, session_tag, Inches(0.5), Inches(0.25), Inches(8), Inches(0.35),
                 font_size=10, color=accent_color, font_name="Calibri")
    add_text_box(slide, title, Inches(0.5), Inches(0.65), Inches(11), Inches(0.5),
                 font_size=20, bold=True, color=WHITE, font_name="Calibri")
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.02), fill_color=accent_color)
    
    # Content blocks
    content_w = Inches(8) if img_path else Inches(12.3)
    block_h = (Inches(6) / len(content_blocks)) - Inches(0.15)
    
    for i, (block_title, block_text, block_color) in enumerate(content_blocks):
        y = Inches(1.35) + i * (block_h + Inches(0.15))
        add_rect(slide, Inches(0.5), y, content_w, block_h, line_color=block_color)
        add_text_box(slide, block_title, Inches(0.65), y + Inches(0.1), content_w - Inches(0.3), Inches(0.35),
                     font_size=13, bold=True, color=WHITE, font_name="Calibri")
        add_text_box(slide, block_text, Inches(0.65), y + Inches(0.5), content_w - Inches(0.3), block_h - Inches(0.65),
                     font_size=12, color=DIM, font_name="Calibri")

def make_activity_slide(prs, session_tag, activity_name, description, tasks, meta_items, accent_color=ACCENT2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    
    add_text_box(slide, session_tag, Inches(0.5), Inches(0.25), Inches(10), Inches(0.35),
                 font_size=10, color=accent_color, font_name="Calibri")
    add_text_box(slide, activity_name, Inches(0.5), Inches(0.65), Inches(11), Inches(0.5),
                 font_size=20, bold=True, color=WHITE, font_name="Calibri")
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.02), fill_color=accent_color)
    
    # Activity description banner
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.0), line_color=accent_color)
    add_text_box(slide, description, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.85),
                 font_size=13, color=DIM, font_name="Calibri")
    
    # Meta chips
    for i, (meta_label, meta_color) in enumerate(meta_items):
        x = Inches(0.5) + i * Inches(3.0)
        add_rect(slide, x, Inches(2.5), Inches(2.8), Inches(0.45), line_color=meta_color)
        add_text_box(slide, meta_label, x + Inches(0.1), Inches(2.55), Inches(2.6), Inches(0.35),
                     font_size=11, color=meta_color, font_name="Calibri", align=PP_ALIGN.CENTER)
    
    # Tasks
    task_y = Inches(3.15)
    for j, (num, task_title, task_text) in enumerate(tasks):
        y = task_y + j * Inches(1.2)
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.05), line_color=ACCENT2)
        add_text_box(slide, f"  {num}  {task_title}", Inches(0.6), y + Inches(0.05), Inches(12), Inches(0.4),
                     font_size=13, bold=True, color=WHITE, font_name="Calibri")
        add_text_box(slide, task_text, Inches(0.8), y + Inches(0.5), Inches(12), Inches(0.5),
                     font_size=12, color=DIM, font_name="Calibri")

def make_end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    
    add_text_box(slide, "CIERRE — UNIDAD 2", Inches(0.5), Inches(0.25), Inches(12), Inches(0.4),
                 font_size=10, color=ACCENT1, font_name="Calibri", align=PP_ALIGN.CENTER)
    add_text_box(slide, "¿Qué logramos en esta Unidad?", Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.8),
                 font_size=32, bold=True, color=WHITE, font_name="Calibri", align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.02), fill_color=ACCENT1)
    
    achievements = [
        "Seleccionar tipos de prueba basándonos en evaluación de riesgo real (Impacto × Probabilidad).",
        "Transformar requisitos monolíticos en Historias de Usuario con metodología INVEST.",
        "Redactar Criterios de Aceptación testeables en sintaxis Gherkin BDD (DADO/CUANDO/ENTONCES).",
        "Diseñar Casos de Prueba formales con los 5 campos estándar.",
        "Implementar Matrices de Trazabilidad (RTM) para auditar cobertura completa.",
        "Construir un Test Plan estratégico maestro con los 3 pilares IEEE 829.",
    ]
    
    cols = 2
    ach_w = Inches(5.9)
    ach_h = Inches(0.85)
    gap_h = Inches(0.12)
    gap_v = Inches(0.12)
    start_x = Inches(0.5)
    start_y = Inches(1.7)
    
    for i, ach in enumerate(achievements):
        col = i % cols
        row = i // cols
        x = start_x + col * (ach_w + Inches(0.3))
        y = start_y + row * (ach_h + gap_v)
        add_rect(slide, x, y, ach_w, ach_h, line_color=ACCENT4)
        add_text_box(slide, "✓  " + ach, x + Inches(0.1), y + Inches(0.08), ach_w - Inches(0.2), ach_h - Inches(0.1),
                     font_size=11, color=DIM, font_name="Calibri")
    
    # Next unit card
    add_rect(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.0), line_color=ACCENT2)
    add_text_box(slide, "PRÓXIMA UNIDAD", Inches(0.7), Inches(5.65), Inches(12), Inches(0.3),
                 font_size=9, color=ACCENT2, font_name="Calibri")
    add_text_box(slide, "Unidad 3 — Ejecución, Automatización y Gestión de Defectos",
                 Inches(0.7), Inches(5.95), Inches(12), Inches(0.25),
                 font_size=13, bold=True, color=WHITE, font_name="Calibri")
    add_text_box(slide, "Bug Tracking con Jira, Selenium, Playwright, métricas de cobertura de código.",
                 Inches(0.7), Inches(6.2), Inches(12), Inches(0.25),
                 font_size=11, color=DIM, font_name="Calibri")


# === MAIN ===
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    print("Generando PPTX — Unidad 2 IF203IINF...")
    
    # Slide 00: Cover
    make_slide_cover(prs)
    print("  [01/18] Portada ✓")
    
    # Slide 01: TOC
    make_toc_slide(prs)
    print("  [02/18] Tabla de Contenidos ✓")
    
    # Slide 02: Section 1
    make_section_slide(prs, "SESIÓN 01",
        "Selección Estratégica de Tipos de Prueba",
        "Cómo decidir qué probar, cuándo y con qué herramientas.",
        ACCENT1,
        [("Principio", "Testing en Contexto (ISTQB)"),
         ("Herramienta", "Matriz Riesgo: Impacto × Probabilidad"),
         ("Práctica", "El Estratega de Calidad")])
    print("  [03/18] Sección 1 ✓")
    
    # Slide 03: Principio del Contexto
    make_content_slide(prs, "SESIÓN 01 · Selección de Tipos",
        "El Principio del Contexto (ISTQB #6)",
        [
            ("Principio 6 ISTQB — El Testing depende del contexto",
             "La selección de tipos de prueba es una decisión de INVERSIÓN FINANCIERA, no solo técnica. Los recursos son limitados: se invierte donde el riesgo es mayor.",
             ACCENT1),
            ("Factores: Complejidad + Regulación + Ciclo de Vida",
             "Sistemas distribuidos, microservicios, requisitos legales (ISO 26262, IEC 62443), y el ciclo Agile vs Waterfall determinan la estrategia de pruebas.",
             ACCENT2),
            ("Ejemplo A — Startup MVP (1 mes para lanzar)",
             "Selección: Pruebas unitarias core, Smoke Tests manuales. Sin pruebas de carga pesadas. El riesgo de negocio es NO llegar al mercado a tiempo.",
             ACCENT4),
        ], ACCENT1)
    print("  [04/18] Principio Contexto ISTQB ✓")
    
    # Slide 04: Matriz de Riesgo
    make_content_slide(prs, "SESIÓN 01 · Selección de Tipos",
        "Matriz de Riesgo: Impacto × Probabilidad",
        [
            ("La Matriz de Riesgo como herramienta de priorización",
             "Clasifica cada funcionalidad según Probabilidad de Fallo × Impacto del Fallo. El cuadrante Crítico recibe cobertura máxima; el cuadrante Bajo puede excluirse.",
             ACCENT1),
            ("CRÍTICO: Alta Probabilidad + Alto Impacto",
             "Pruebas exhaustivas: unitarias + integración + E2E + seguridad. Ejecución en cada build. Ejemplo: módulo de pagos en e-commerce de alto volumen.",
             ERR),
            ("BAJO: Baja Probabilidad + Bajo Impacto",
             "Se documenta la decisión de NO probar. El riesgo es asumido explícitamente por el Product Owner. Ejemplo: página estática 'Acerca de nosotros'.",
             ACCENT4),
        ], ACCENT1, img_path=str(IMAGES_DIR / "risk_matrix.png"))
    print("  [05/18] Matriz de Riesgo ✓")
    
    # Slide 05: Activity Estratega
    make_activity_slide(prs, "ACTIVIDAD PRÁCTICA — SESIÓN 01", "El Estratega de Calidad",
        "Asume el rol de Test Manager. Se presentan 3 proyectos ficticios. Elige EXACTAMENTE 3 tipos de prueba para cada uno y justifica con Matriz de Riesgo.",
        [
            ("01", "E-Commerce — Black Friday", "Miles de usuarios simultáneos. ¿Qué 3 tipos de prueba priorizas? Justifica con Impacto × Probabilidad."),
            ("02", "Control de Semáforos IoT", "Sistema gubernamental 24/7. ¿Qué normas regulatorias aplican? ¿Cómo justificas excluir pruebas?"),
            ("03", "App de Recetas (Red Social)", "Equipo de 2 devs, presupuesto mínimo. ¿Cuál es el nivel MÍNIMO aceptable de testing?"),
        ],
        [("📤 Buzón EVA", ACCENT2), ("⏱ 30 minutos", ACCENT3), ("👥 Grupal (2-3)", ACCENT4)],
        ACCENT2)
    print("  [06/18] Actividad Estratega ✓")
    
    # Slide 06: Section 2
    make_section_slide(prs, "SESIÓN 02",
        "Requerimientos e Historias de Usuario",
        "De la ambigüedad del cliente a la especificidad técnica.",
        ACCENT2,
        [("Comparativa", "SRS (Waterfall) vs User Stories"),
         ("Herramienta", "Criterio INVEST + Gherkin BDD"),
         ("Práctica", "El Traductor Ágil")])
    print("  [07/18] Sección 2 ✓")
    
    # Slide 07: SRS vs User Stories
    make_content_slide(prs, "SESIÓN 02 · Requerimientos",
        "SRS vs Historia de Usuario — Waterfall vs Agile",
        [
            ("Requerimiento Tradicional SRS (Waterfall)",
             'Documentos extensos, enfocados en el "Qué" desde una perspectiva de ingeniería. "El sistema deberá permitir autenticación contra LDAP con SHA-256 en < 200ms."',
             ACCENT2),
            ("Historia de Usuario (Agile) — COMO / QUIERO / PARA",
             '"COMO empleado remoto, QUIERO ingresar con mi clave corporativa, PARA no recordar contraseñas distintas." Un recordatorio para tener una conversación.',
             ACCENT1),
            ("Criterio INVEST — La Historia de Usuario debe ser:",
             "I=Independiente  N=Negociable  V=Valiosa  E=Estimable  S=Small (Pequeña)  T=Testeable ★ (La T es la que conecta con Gherkin)",
             ACCENT4),
        ], ACCENT2, img_path=str(IMAGES_DIR / "user_stories.png"))
    print("  [08/18] SRS vs User Stories ✓")
    
    # Slide 08: INVEST + Gherkin
    make_content_slide(prs, "SESIÓN 02 · Requerimientos",
        "Criterios de Aceptación — INVEST + Gherkin BDD",
        [
            ("Sintaxis Gherkin — El estándar de facto para BDD",
             "Permite que Negocio, Desarrollo y QA hablen el mismo idioma. Directamente ejecutable con Cucumber, Behave, SpecFlow. Escenario: DADO/CUANDO/ENTONCES.",
             ACCENT2),
            ("Ejemplo — Suscripción Newsletter (3 escenarios)",
             "ÉXITO: DADO que el correo no existe, CUANDO se suscribe, ENTONCES se confirma.  |  DUPLICADO: DADO que existe, ENTONCES avisa.  |  INVÁLIDO: ENTONCES error de formato.",
             ACCENT1),
            ("¿Por qué Gherkin sobre criterios en texto libre?",
             "Es directamente automatizable. Elimina interpretaciones. Sirve como documentación viva que siempre refleja el comportamiento real del sistema en producción.",
             ACCENT4),
        ], ACCENT2)
    print("  [09/18] INVEST + Gherkin ✓")
    
    # Slide 09: Activity Traductor
    make_activity_slide(prs, "ACTIVIDAD PRÁCTICA — SESIÓN 02", "El Traductor Ágil",
        "El Gerente de Ventas envió un correo con un requerimiento monolítico. Debes desglosarlo al formato Ágil profesional.",
        [
            ("01", "Redacta la Historia de Usuario", "COMO / QUIERO / PARA capturando el valor de negocio real."),
            ("02", "Redacta 3 Criterios de Aceptación Gherkin", "Cubre: éxito (correo válido), duplicado y formato inválido (sin @)."),
            ("03", "Verifica con INVEST", "¿Tu historia cumple todos los criterios? Especialmente: ¿es Testeable?"),
        ],
        [("💬 Foro EVA", ACCENT2), ("⏱ 45 minutos", ACCENT3), ("👤 Individual", ACCENT1)],
        ACCENT1)
    print("  [10/18] Actividad Traductor ✓")
    
    # Slide 10: Section 3
    make_section_slide(prs, "SESIÓN 03",
        "Diseño de Casos de Prueba y Trazabilidad",
        "La anatomía de un Test Case infalible. Cobertura exhaustiva con RTM.",
        ACCENT3,
        [("Estructura", "Anatomía del Test Case (5 campos)"),
         ("Herramienta", "RTM — Requirement Traceability Matrix"),
         ("Práctica", "Disección del Bug")])
    print("  [11/18] Sección 3 ✓")
    
    # Slide 11: Anatomía TC
    make_content_slide(prs, "SESIÓN 03 · Test Cases",
        "Anatomía del Caso de Prueba (Test Case)",
        [
            ("Campo 01: ID / Título — Identificador único y nombre descriptivo",
             "TC-001: Validar login exitoso con credenciales correctas — Módulo de autenticación. Permite referenciar el caso en RTM, defectos y reportes de calidad.",
             ACCENT3),
            ("Campos 02-03: Precondiciones y Pasos",
             "PRECOND: Estado exacto del sistema ANTES de ejecutar. SIN precondiciones, el resultado no es reproducible. PASOS: Acciones explícitas y atómicas. Cada paso = una acción.",
             ACCENT1),
            ("Campos 04-05: Resultado Esperado y Datos de Prueba",
             "RESULTADO: La validación que define el ÉXITO. El sistema redirige a /dashboard y muestra 'Bienvenido Admin'. DATOS: user='admin', pass='1234', env=staging-v2.",
             ACCENT4),
        ], ACCENT3, img_path=str(IMAGES_DIR / "test_case.png"))
    print("  [12/18] Anatomía Test Case ✓")
    
    # Slide 12: RTM
    make_content_slide(prs, "SESIÓN 03 · Test Cases",
        "Requirement Traceability Matrix (RTM)",
        [
            ("¿Qué es la RTM? — El mapa de cobertura del proyecto",
             "Documento (tabla o gestionado en Jira+Xray) que mapea y rastrea CADA requerimiento con sus Casos de Prueba respectivos. El contrato de calidad del proyecto.",
             ACCENT3),
            ("Objetivos: Cobertura, Impacto y Auditoría",
             "100% Cobertura: ningún requisito sin probar. Análisis de Impacto: si un req cambia, sé qué TCs actualizar. Auditoría: demuestra al cliente que el sistema fue validado.",
             ACCENT1),
            ("En la tabla: REQ-03 Bloqueo tras 3 intentos — ¡FALTA TEST!",
             "La RTM detecta automáticamente brechas de cobertura. REQ-03 tiene 0% de cobertura. Es una vulnerabilidad de seguridad sin caso de prueba asignado. Acción requerida.",
             ERR),
        ], ACCENT3)
    print("  [13/18] RTM ✓")
    
    # Slide 13: Activity Bug
    make_activity_slide(prs, "ACTIVIDAD PRÁCTICA — SESIÓN 03", "Disección del Bug",
        "Diseña un Caso de Prueba estructurado que exponga un defecto crítico reportado en producción. Recrea el bug de forma determinista.",
        [
            ("01", "Título del TC", "Nombre que incluya 'doble click' y 'sin fondos'. Identifica el bug sin verlo ejecutar."),
            ("02", "Precondiciones + Pasos precisos", "Incluir el TIMING del doble click. ¿En qué momento exacto debe ser el segundo click?"),
            ("03", "Resultado Esperado (el correcto)", "Inventario sin cambios. Mensaje de error de pago. NO el comportamiento bugueado actual."),
        ],
        [("📤 Buzón EVA", ACCENT3), ("⏱ 40 minutos", ERR), ("👤 Individual", ACCENT1)],
        ERR)
    print("  [14/18] Actividad Bug ✓")
    
    # Slide 14: Section 4
    make_section_slide(prs, "SESIÓN 04 — SINCRÓNICA",
        "El Plan de Pruebas Maestro y Dinámicas Finales",
        "Unificar requerimientos, estrategias, casos y recursos en el documento rector.",
        ACCENT4,
        [("Documento", "Test Plan IEEE 829 — 3 Pilares"),
         ("Roleplay", "Defendiendo el Plan (20 min)"),
         ("Dinámica", "QA Audit — Auditoría Cruzada")])
    print("  [15/18] Sección 4 ✓")
    
    # Slide 15: IEEE 829
    make_content_slide(prs, "SESIÓN 04 · Test Plan",
        "IEEE 829: Los 3 Pilares del Plan de Pruebas",
        [
            ("Pilar 1 — Alcance (In / Out Scope)",
             "Define explícitamente QUÉ se va a probar y qué NO. Aprobado por el Product Owner. Acota la responsabilidad del QA y previene el scope creep.",
             ACCENT4),
            ("Pilar 2 — Estrategia y Entorno",
             "Define CÓMO se probará. Tipos de prueba seleccionados, herramientas (Jira+Xray, Selenium, Postman, k6), topología del entorno staging, y cronograma.",
             ACCENT1),
            ("Pilar 3 — Criterios y Entregables",
             "ENTRADA: Código compilado, ambiente disponible, RTM lista. SALIDA: 100% bugs críticos cerrados, RTM ≥ 95% PASS. Sin criterios de salida, las pruebas nunca terminan.",
             ACCENT3),
        ], ACCENT4, img_path=str(IMAGES_DIR / "test_plan.png"))
    print("  [16/18] IEEE 829 Test Plan ✓")
    
    # Slide 16: Dinámicas
    make_content_slide(prs, "ACTIVIDADES SINCRÓNICAS — SESIÓN 04",
        "Roleplay + QA Audit en Vivo",
        [
            ("Defendiendo el Plan — Presentación Ejecutiva (20 min)",
             "Los grupos presentan su Test Plan ante el profesor (rol de Gerente exigente). Deben defender por qué excluyeron pruebas del alcance para cumplir el deadline.",
             ACCENT4),
            ("QA Audit — Auditoría Cruzada (Breakout Rooms, 40 min)",
             "Grupo A intercambia Test Cases con Grupo B. Cada grupo busca fisuras: pasos incompletos, precondiciones imposibles, resultados no verificables. Se premia al mejor auditor.",
             RGBColor(0xf5, 0x9e, 0x0b)),
        ], ACCENT4)
    print("  [17/18] Dinámicas Sincrónicas ✓")
    
    # Slide 17: Eval Context
    make_content_slide(prs, "EVALUACIÓN SUMATIVA 2 (40%)",
        "Continuidad del Repositorio CS50 / LeetCode",
        [
            ("Retomando el Código",
             "En la Evaluación 1 definiste el tipo de prueba y diseñaste 4 casos aislados para el repositorio.",
             ERR),
            ("Evolución del Rol — QA Lead",
             "Tu objetivo ahora es diseñar la estrategia de cobertura total (Test Plan) y asegurar que los enunciados de los algoritmos se traduzcan en criterios automatizables BDD.",
             RGBColor(0xf5, 0x9e, 0x0b)),
        ], ERR)
    print("  [18/21] Eval 2 Contexto ✓")
    
    # Slide 18: Eval Parts
    make_content_slide(prs, "EVALUACIÓN SUMATIVA 2 (40%)",
        "Partes de la Evaluación",
        [
            ("Parte 1: Matriz de Pruebas (RTM) - 35 pts",
             "Hoja de cálculo. Mapea todos los requerimientos algorítmicos con sus casos de prueba positivos y negativos. Cobertura total.",
             ACCENT1),
            ("Parte 2: Diseño del Plan y Casos - 35 pts",
             "Hoja de cálculo. Cada caso debe tener ID, Precondiciones, Pasos, Datos y Resultado Esperado.",
             ACCENT2),
            ("Parte 3: Criterios Gherkin - 30 pts",
             "PDF (Max 10 pág). Traducir requerimientos usando DADO, CUANDO, ENTONCES y AND. Vocabulario de negocio.",
             RGBColor(0xf5, 0x9e, 0x0b)),
        ], RGBColor(0xf5, 0x9e, 0x0b))
    print("  [19/21] Eval 2 Partes ✓")
    
    # Slide 19: Eval Rubric
    make_content_slide(prs, "CRITERIOS DE CALIFICACIÓN",
        "Rúbrica y Restricciones Formales",
        [
            ("Restricción: Formato y Entrega",
             "Subir PDF y .xls directo a plataforma (no ZIP/RAR). Portada formal con RUT de todos los integrantes (nota 1.0 si un alumno no aparece).",
             ERR),
            ("Restricción: Uso de Gherkin",
             "Se auditará coherencia técnica. Prohibido usar Gherkin con lenguaje de código interno (ej: el objeto JSON tiene id=1).",
             ERR),
        ], ERR)
    print("  [20/21] Eval 2 Rubrica ✓")

    # Slide 20: End
    make_end_slide(prs)
    print("  [21/21] Cierre Unidad 2 ✓")
    
    # Save
    prs.save(str(OUTPUT_PPTX))
    print(f"\n✅ PPTX generado exitosamente: {OUTPUT_PPTX}")
    print(f"   Slides: 21  |  Tamaño: {OUTPUT_PPTX.stat().st_size // 1024} KB")

if __name__ == "__main__":
    main()
